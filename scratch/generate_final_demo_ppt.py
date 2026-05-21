"""
Generate HPE CPP-3 Final Demo PowerPoint.
Sources: create-ppt.txt (mentor brief + team deck + user guide outline)
Figures 1-6: screenshot placeholders only (insert images manually in PowerPoint).
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "HPE_CPP3_Final_Demo.pptx"
IMAGES = ROOT / "images"

# HPE-inspired palette
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
HPE_GREEN = RGBColor(0x01, 0xA9, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
MID_GRAY = RGBColor(0xD9, 0xDE, 0xE3)
TEXT = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
PLACEHOLDER_BG = RGBColor(0xE9, 0xEE, 0xF5)

FIGURES = [
    {
        "num": 1,
        "title": "Authentication & Sign-Up",
        "file": "Sign-up.png",
        "caption": "Role-based registration (team member, lead, admin) — MongoDB + JWT via chatbot service",
        "bullets": [
            "React LoginPage.jsx · chatbot POST /api/auth/register",
            "Team/cluster assignment for test ring visibility",
        ],
    },
    {
        "num": 2,
        "title": "Discovery Tab — Live BFS Scan",
        "file": "discovery-tab-fullscreen.png",
        "caption": "Automated topology crawl with SSE event stream into Neo4j and Elasticsearch",
        "bullets": [
            "POST /api/discovery/start · GET /api/discovery/stream",
            "Simulator port 5001 · DiscoveryPanel + TopologyCanvas",
        ],
    },
    {
        "num": 3,
        "title": "Test Ring Viewer — Interactive Topology",
        "file": "test-ring-viewer-tab.png",
        "caption": "SAN diagram and graph canvas for arrays, switches, and hosts",
        "bullets": [
            "SANDiagram.jsx + TopologyCanvas.jsx · /api/graph/neo4j",
            "Role-based team filters from teamconfig.json",
        ],
    },
    {
        "num": 4,
        "title": "Inventory — Hierarchical Resources",
        "file": "inventory-tab.png",
        "caption": "Path-oriented tree from ArraySystem down to PhysicalDisk",
        "bullets": [
            "HierarchyTree.jsx with search and path highlighting",
            "Read-only graph consumption from Neo4j",
        ],
    },
    {
        "num": 5,
        "title": "AI Assistant — SAN Agent & GraphRAG",
        "file": "chatbot-san-agent-demo-query.png",
        "caption": "Natural-language queries with Neo4j context, Groq/Gemini, and agent plans",
        "bullets": [
            "Also: neo4j-powered-chatbot.png · radial-quick-query-menu.png",
            "Chatbot :5010 · Flask RAG /api/chat · AgentStepTimeline",
        ],
    },
    {
        "num": 6,
        "title": "Health Tab — Service & Infrastructure Status",
        "file": "list-of-services-from-health-tab.png",
        "caption": "Flask API, Neo4j, Elasticsearch, chatbot, and simulator connectivity",
        "bullets": [
            "GET /api/health · Investigate → pre-fills AI Assistant",
            "Capacity bars, issue list, recommendations",
        ],
    },
]


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)


def add_bullets(slide, items, left=0.55, top=1.25, width=12.2, height=5.8, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)


def slide_title(prs, title, subtitle, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    tbox = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(1.5))
    p = tbox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    sbox = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12), Inches(1.2))
    p2 = sbox.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = HPE_GREEN

    if footer:
        fbox = slide.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(12), Inches(0.8))
        p3 = fbox.text_frame.paragraphs[0]
        p3.text = footer
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    return slide


def slide_content(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, title, subtitle)
    add_bullets(slide, bullets, top=1.35)
    return slide


def slide_two_column(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, title)
    lt = slide.shapes.add_textbox(Inches(0.55), Inches(1.2), Inches(6.1), Inches(0.45))
    lt.text_frame.paragraphs[0].text = left_title
    lt.text_frame.paragraphs[0].font.bold = True
    lt.text_frame.paragraphs[0].font.size = Pt(16)
    lt.text_frame.paragraphs[0].font.color.rgb = HPE_GREEN
    add_bullets(slide, left_items, left=0.55, top=1.65, width=6.0, size=15)

    rt = slide.shapes.add_textbox(Inches(6.85), Inches(1.2), Inches(6.1), Inches(0.45))
    rt.text_frame.paragraphs[0].text = right_title
    rt.text_frame.paragraphs[0].font.bold = True
    rt.text_frame.paragraphs[0].font.size = Pt(16)
    rt.text_frame.paragraphs[0].font.color.rgb = HPE_GREEN
    add_bullets(slide, right_items, left=6.85, top=1.65, width=6.0, size=15)
    return slide


def slide_figure_placeholder(prs, fig):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(
        slide,
        f"Figure {fig['num']}: {fig['title']}",
        "Live demo screenshot — insert image in placeholder below",
    )

    # Placeholder box
    left, top, w, h = Inches(0.55), Inches(1.35), Inches(12.2), Inches(4.85)
    ph = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, w, h)
    ph.fill.solid()
    ph.fill.fore_color.rgb = PLACEHOLDER_BG
    ph.line.color.rgb = MID_GRAY
    ph.line.width = Pt(1.5)

    tf = ph.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ INSERT FIGURE {fig['num']} ]\n\n"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run2 = p.add_run()
    run2.text = f"images/{fig['file']}\n\n{fig['caption']}"
    run2.font.size = Pt(14)
    run2.font.color.rgb = MUTED

    add_bullets(slide, fig["bullets"], top=6.35, height=1.0, size=13)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title ---
    slide_title(
        prs,
        "HPE CPP-3 — Final Demo",
        "Topology Intelligence & Health Chatbot\nfor Storage Test Ring Environments",
        "Team: Preetham M R · Sachethan I · Samarth K N · Unmesh Raj  |  Mentor: Dr. Mohana",
    )

    slide_content(
        prs,
        "Agenda",
        [
            "Mentor problem statement & deliverables",
            "Our solution and system architecture (monorepo)",
            "Core capabilities: discovery, graph, simulator, RAG chatbot",
            "Live dashboard demo — Figures 1–6",
            "Outcomes and future scope",
        ],
    )

    # --- Mentor brief (source 1) ---
    slide_content(
        prs,
        "Mentor Problem Statement",
        [
            "View all equipment connectivity and component health in one place",
            "Maintain connectivity details without fragmented spreadsheets",
            "Infrastructure Analytics Engine: query FC, iSCSI, NVMe/TCP support and applicable test cases",
            "Capture Array / Host / Infrastructure → store in repository → query stored details",
        ],
        subtitle="Hardware Test Ring Management Tool (HPE Alletra MP B10000 test rings)",
    )

    slide_two_column(
        prs,
        "Mentor — Example Query Outcomes",
        "Host queries",
        [
            "Given a host: arrays zoned with, HBA/driver/FW",
            "Given a host: connected switch",
        ],
        "Array queries",
        [
            "Hosts zoned with + OS type; protocols supported",
            "Switch state, TPD version, cage/PD health, failed PDs",
            "Arrays >200TB, NVMe/TCP tests, online upgrade capability",
        ],
    )

    slide_content(
        prs,
        "Mentor — Implementation Breakdown",
        [
            "(a) Extract key info from arrays, hosts, switches via CLI on periodic/need basis → DB",
            "(b) Extensible DB schema + export to spreadsheet",
            "(c) Framework: English statements → database/graph query",
            "(d) Web UI for stored details and query results",
            "→ Our delivery: Neo4j graph + Elasticsearch + React dashboard + SAN agent chatbot",
        ],
    )

    # --- Problem / Solution (sources 2 & 3) ---
    slide_content(
        prs,
        "The Problem: Complexity Without Intelligence",
        [
            "Hundreds of hosts, switches, HBAs, controllers, arrays — knowledge in spreadsheets & tribal expertise",
            "No single source of truth; manual connectivity tracing",
            "Reactive health visibility; opaque zoning relationships",
            "Multiple unstructured logs slow troubleshooting",
        ],
    )

    slide_two_column(
        prs,
        "Our Solution — Platform Pillars",
        "Data & graph",
        [
            "Neo4j: ArraySystem, Host, Switch, Cage, Node, PhysicalDisk",
            "Elasticsearch: discovery metadata for RAG",
            "MongoDB: users, chat history",
        ],
        "Intelligence & UI",
        [
            "Python SAN simulator (digital twin) — port 5001",
            "BFS discovery engine with live SSE streaming",
            "GraphRAG + SAN agent (Groq/Gemini) — port 5010",
            "React 18 dashboard — Discovery, Topology, Inventory, Emulator, AI, Health",
        ],
    )

    # --- Architecture (codebase-accurate) ---
    slide_content(
        prs,
        "System Architecture (Implemented)",
        [
            "Dashboard (Vite :5173) → Flask API (:5005) → Simulator (:5001)",
            "Flask → Neo4j (:7687) + Elasticsearch (:9200)",
            "Dashboard → Chatbot (:5010) → MongoDB (:27017) + LLM APIs",
            "Discovery: BFS crawler → fingerprint → parse → neo4j_stored / es_indexed events",
            "No production hardware required for demo — full virtual test ring",
        ],
        subtitle="Monorepo — docker-compose + npm start",
    )

    slide_two_column(
        prs,
        "Query Engine — Intent to Graph",
        "How it works",
        [
            "Natural language → intent classification (agent / GraphRAG / standard)",
            "Elasticsearch retrieval + Neo4j Cypher traversal",
            "SAN agent: multi-step plans + optional simulator CLI execution",
            "Deterministic graph answers — not generic hallucination-only chat",
        ],
        "Example queries (mentor-aligned)",
        [
            "Given array PROD-A: hosts zoned with + OS type",
            "List arrays with failed PDs; TPD version for PROD-A",
            "Arrays with >200TB usable space; switch state",
            "Radial menu: one-click prompts in dashboard",
        ],
    )

    slide_content(
        prs,
        "Digital Twin — SAN Simulator",
        [
            "simulator_manager.py boots virtual arrays, FC switches, Linux/Windows hosts",
            "HPE CLI replay (showsys, showport, showpd) + host commands (multipath -ll)",
            "REST /api/simulator/exec — used by Discovery, Emulator, and AI agent",
            "Emulator tab: device picker + quick commands + live terminal",
            "Figure 4 (Inventory) and Emulator terminal shown in Figures 3–5 demo slides",
        ],
    )

    slide_content(
        prs,
        "Technology Stack",
        [
            "Frontend: React 18 + Vite",
            "API: Flask + Python 3.10 (api/, discovery/)",
            "Chatbot: Node.js Express (chatbot-service/)",
            "Graph: Neo4j 5.18 | Search: Elasticsearch 8.13 | Docs: MongoDB",
            "LLM: Groq + Gemini | Infra: Docker Compose",
            "Quick start: docker-compose up -d neo4j elasticsearch mongo && npm start",
        ],
    )

    # --- Demo section divider ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, HPE_GREEN)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(12), Inches(2))
    p = box.text_frame.paragraphs[0]
    p.text = "Live Demo — Dashboard Walkthrough"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = box.text_frame.add_paragraph()
    p2.text = "Figures 1–6: insert screenshots from images/ folder"
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    # Optional: emulator as callout before fig 5 — user has 6 figs mapped above
    slide_figure_placeholder(prs, FIGURES[0])
    slide_figure_placeholder(prs, FIGURES[1])
    slide_figure_placeholder(prs, FIGURES[2])
    slide_figure_placeholder(prs, FIGURES[3])
    slide_figure_placeholder(prs, FIGURES[4])
    slide_figure_placeholder(prs, FIGURES[5])

    # Extra note slide for auxiliary screenshots
    slide_content(
        prs,
        "Additional UI Assets (optional slides)",
        [
            "terminal-simulator-tab.png — Emulator / digital twin CLI (Fig 5 alternate)",
            "neo4j-powered-chatbot.png — Graph-enriched answer detail",
            "radial-quick-query-menu.png — Radial quick-query menu (User Guide § UI)",
            "Duplicate or split Figures 5–6 if your demo script needs separate chat vs health timing",
        ],
    )

    # --- Outcomes ---
    slide_two_column(
        prs,
        "Delivered vs Mentor Requirements",
        "Delivered",
        [
            "✓ Connectivity & topology visualization (Discovery + Test Ring Viewer)",
            "✓ Component health (Health tab + graph node status)",
            "✓ Repository: Neo4j + Elasticsearch + MongoDB",
            "✓ English → query: SAN agent + GraphRAG chatbot",
            "✓ Web UI: full React dashboard with role-based access",
        ],
        "Future scope",
        [
            "Live hardware SNMP/SSH collectors (beyond simulator)",
            "Periodic scheduled discovery jobs",
            "Predictive analytics & automated remediation",
            "Kubernetes / multi-cluster deployment",
        ],
    )

    slide_content(
        prs,
        "Expected Outcomes",
        [
            "Single source of truth for SAN test ring topology",
            "Faster troubleshooting: path trace + AI Investigate workflow",
            "Proactive health: service badges, capacity, failed/degraded nodes",
            "Safe experimentation via digital twin without physical arrays",
            "Scalable monorepo architecture for HPE lab environments",
        ],
    )

    slide_title(prs, "Thank You", "Questions & live demo", "HPE Ring Test Management — monorepo")

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
